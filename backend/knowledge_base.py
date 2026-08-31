"""Read-only local knowledge base for the user's mathematical-modelling pack.

The adapter deliberately does not copy, execute, or mutate source material.  It
keeps a small in-process index for the prototype; a production deployment can
persist the same records in SQLite FTS5 without changing the public methods.
"""
from __future__ import annotations

import hashlib
import html
import os
import re
import threading
import time
import zipfile
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional

DEFAULT_ROOT = r"C:\Users\zyy20\Desktop\数学建模资料全套包"
TEMP_SUFFIXES = {".qkdownloading", ".part", ".crdownload", ".tmp"}
ALLOWED_EXTENSIONS = {
    ".pdf", ".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx", ".csv",
    ".txt", ".md", ".html", ".htm", ".json", ".yaml", ".yml", ".m", ".py",
    ".r", ".rmd", ".tex", ".xml", ".log", ".h", ".hpp", ".c", ".cc", ".cpp",
    ".java", ".js", ".ts", ".sql", ".query", ".ltx", ".mat", ".xyz", ".lg4",
    ".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".svg",
}
TEXT_EXTENSIONS = {
    ".txt", ".md", ".html", ".htm", ".csv", ".json", ".yaml", ".yml",
    ".m", ".py", ".r", ".rmd", ".tex", ".xml", ".log", ".h", ".hpp", ".c",
    ".cc", ".cpp", ".java", ".js", ".ts", ".sql", ".query", ".ltx",
}
# Formats for which the bounded adapter can make a first-pass body/text
# extraction.  This is deliberately narrower than ``ALLOWED_EXTENSIONS``:
# images and legacy Office binaries can be opened by the Owner but are not
# treated as searchable body text.
OPENXML_EXTENSIONS = {".docx", ".pptx", ".xlsx"}
BODY_EXTRACTABLE_EXTENSIONS = set(TEXT_EXTENSIONS) | {".pdf"} | OPENXML_EXTENSIONS
# Old binary Office and macro/container variants are retained as metadata
# records.  We expose them separately so a caller does not mistake an
# indexed filename for understood document content.
LEGACY_OFFICE_EXTENSIONS = {
    ".doc", ".dot", ".ppt", ".pps", ".xls", ".xlt", ".wps", ".et", ".dps",
    ".docm", ".dotm", ".pptm", ".ppsm", ".xlsm", ".xltm",
}
KIND_BY_EXTENSION = {
    ".pdf": "pdf", ".doc": "document", ".docx": "document", ".ppt": "presentation",
    ".pptx": "presentation", ".xls": "spreadsheet", ".xlsx": "spreadsheet",
    ".csv": "spreadsheet", ".m": "code", ".py": "code", ".r": "code",
    ".rmd": "code", ".tex": "code", ".html": "html", ".htm": "html",
    ".md": "text", ".txt": "text", ".json": "text", ".yaml": "text",
    ".yml": "text", ".xml": "text", ".log": "text", ".h": "code", ".hpp": "code",
    ".c": "code", ".cc": "code", ".cpp": "code", ".java": "code", ".js": "code",
    ".ts": "code", ".sql": "code", ".query": "code", ".ltx": "code",
    ".png": "asset", ".jpg": "asset", ".jpeg": "asset", ".bmp": "asset",
    ".gif": "asset", ".webp": "asset", ".svg": "asset",
}
DOC_ID_RE = re.compile(r"^kbdoc_[0-9a-f]{16}$")
TOKEN_RE = re.compile(r"[A-Za-z0-9_]+|[\u4e00-\u9fff]")
YEAR_RE = re.compile(r"(?:19|20)\d{2}")
# Keep inventory responsive for multi-gigabyte packs.  Selected documents are
# content-hashed on demand; callers may raise this limit explicitly for a
# smaller corpus with ``GAOJIAO_KB_INLINE_HASH_LIMIT``.
try:
    INLINE_HASH_LIMIT = max(0, int(os.getenv("GAOJIAO_KB_INLINE_HASH_LIMIT", "0")))
except ValueError:
    INLINE_HASH_LIMIT = 0
INLINE_TEXT_LIMIT = 2 * 1024 * 1024
OFFICE_SOURCE_LIMIT = 16 * 1024 * 1024
OFFICE_UNPACK_LIMIT = 64 * 1024 * 1024
try:
    MAX_OPEN_FILE_BYTES = max(1, int(os.getenv("GAOJIAO_KB_MAX_OPEN_BYTES", str(256 * 1024 * 1024))))
except ValueError:
    MAX_OPEN_FILE_BYTES = 256 * 1024 * 1024
try:
    ON_DEMAND_HASH_LIMIT = max(0, int(os.getenv("GAOJIAO_KB_ON_DEMAND_HASH_LIMIT", str(64 * 1024 * 1024))))
except ValueError:
    ON_DEMAND_HASH_LIMIT = 64 * 1024 * 1024
try:
    KB_CACHE_TTL_SECONDS = max(0.0, float(os.getenv("GAOJIAO_KB_CACHE_TTL_SECONDS", "60")))
except ValueError:
    KB_CACHE_TTL_SECONDS = 60.0


def _now() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds")


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for block in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def _decode(data: bytes) -> str:
    for encoding in ("utf-8-sig", "gb18030", "big5", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _strip_html(value: str) -> str:
    value = re.sub(r"<script[\s\S]*?</script>|<style[\s\S]*?</style>", " ", value, flags=re.I)
    value = re.sub(r"<[^>]+>", " ", value)
    return re.sub(r"\s+", " ", html.unescape(value)).strip()


def _tokens(value: str) -> List[str]:
    raw = TOKEN_RE.findall(value)
    chinese_count = len(re.findall(r"[\u4e00-\u9fff]", value))
    base = [item.lower() for item in raw if not re.fullmatch(r"[\u4e00-\u9fff]", item) or chinese_count <= 1]
    # Chinese bigrams make a query such as “敏感性” useful without requiring a
    # heavyweight tokenizer.  For multi-character Chinese queries, do not add
    # every single character: common characters such as “的/模/型” otherwise
    # make almost the entire corpus look like a metadata hit.  A one-character
    # query is retained as an intentional rare-term lookup.
    compact = re.sub(r"\s+", "", value)
    chinese_count = len(re.findall(r"[\u4e00-\u9fff]", compact))
    segments = re.findall(r"[A-Za-z0-9_]+|[\u4e00-\u9fff]+", value)
    for segment in segments:
        if not re.search(r"[\u4e00-\u9fff]", segment):
            continue
        if len(segment) == 1:
            if chinese_count <= 1:
                base.append(segment.lower())
        else:
            base.extend(segment[i : i + 2].lower() for i in range(len(segment) - 1))
    # Preserve the common mixed notation “B题/C题” without creating bigrams
    # across unrelated whitespace-separated words.
    base.extend(item.replace(" ", "").lower() for item in re.findall(r"[A-Za-z]\s*[\u4e00-\u9fff]", value))
    return list(dict.fromkeys(base))


class KnowledgeBase:
    """Read-only inventory/search facade.

    ``root`` is an explicit test override.  In normal use the environment
    variable is preferred and then the documented Windows default is used.
    """

    def __init__(self, root: Optional[str] = None):
        configured = root or os.getenv("GAOJIAO_MATERIALS_ROOT") or DEFAULT_ROOT
        self.root = Path(configured).expanduser()
        self._root_resolved: Optional[Path] = None
        self._records: Dict[str, Dict[str, Any]] = {}
        self._summary_cache: Optional[Dict[str, Any]] = None
        self._summary_cached_at = 0.0
        self._text_cache: Dict[str, str] = {}
        self._text_cache_pages: Dict[str, int] = {}
        self._lock = threading.RLock()

    @property
    def root_id(self) -> str:
        return "math-modeling-pack"

    def _root(self) -> Optional[Path]:
        if not self.root.exists() or not self.root.is_dir():
            return None
        try:
            resolved = self.root.resolve(strict=True)
        except OSError:
            return None
        self._root_resolved = resolved
        return resolved

    def _relative(self, path: Path, root: Path) -> Optional[str]:
        try:
            resolved = path.resolve(strict=True)
            resolved.relative_to(root)
        except (OSError, ValueError):
            return None
        return resolved.relative_to(root).as_posix()

    def _classify(self, rel: str, path: Path, size: int, digest: Optional[str], mtime_ns: int,
                  ctime_ns: int = 0, inode: int = 0) -> Dict[str, Any]:
        parts = rel.split("/")
        module = parts[0] if parts else "未分类"
        extension = path.suffix.lower()
        lower_rel = rel.lower()
        # Directory semantics are stronger than a bare extension for this
        # corpus: a PDF under “官方评阅标准” is a rubric, while a PDF under
        # “获奖论文” is a paper.  These labels are retrieval facets only.
        if any(word in lower_rel for word in ("评阅", "评审标准", "官方答案", "评阅要点")):
            kind = "rubric"
        elif any(word in lower_rel for word in ("获奖论文", "优秀论文", "论文合集", "论文精选")):
            kind = "paper"
        elif any(word in lower_rel for word in ("赛题解析", "历年赛题", "真题", "赛题", "竞赛通知")):
            kind = "problem"
        elif any(word in lower_rel for word in ("模板", "写作", "论文格式", "排版", "备赛")):
            kind = "template"
        elif any(word in lower_rel for word in ("代码", "算法")) or extension in {".m", ".py", ".r", ".rmd", ".tex", ".h", ".hpp", ".c", ".cc", ".cpp", ".java", ".js", ".ts", ".sql", ".query", ".ltx"}:
            kind = "code"
        elif any(word in lower_rel for word in ("课程", "课件", "教材", "学习资料")):
            kind = "course"
        elif any(word in lower_rel for word in ("软件", "工具", "指南")):
            kind = "tool"
        elif any(word in lower_rel for word in ("参考书", "书库")):
            kind = "book"
        else:
            kind = KIND_BY_EXTENSION.get(extension, "other")
        current_year = datetime.now(timezone.utc).year
        years = sorted({int(value) for value in YEAR_RE.findall(rel) if 1990 <= int(value) <= current_year + 1})
        tags = [f"module:{module}", f"extension:{extension.lstrip('.') or 'none'}", f"kind:{kind}"]
        tags.extend(f"year:{year}" for year in years)
        # Common A/B/C problem-code convention; this is a search tag only.
        codes = re.findall(r"(?<![A-Za-z])([A-E])\d{2,4}(?!\d)", path.stem, flags=re.I)
        tags.extend(f"problem:{code.upper()}" for code in codes)
        doc_id = "kbdoc_" + hashlib.sha256(rel.encode("utf-8")).hexdigest()[:16]
        return {
            "doc_id": doc_id, "path_rel": rel, "title": path.stem,
            "root_id": self.root_id, "content_class": "user_private",
            "module": module, "extension": extension, "kind": kind, "years": years,
            "tags": list(dict.fromkeys(tags)), "size": size, "size_bytes": size,
            "sha256": ("sha256:" + digest) if digest else None,
            "hash_status": "INLINE" if digest else "DEFERRED",
            "mtime_ns": mtime_ns, "ctime_ns": ctime_ns, "inode": inode,
            "source_status": "LOCAL_INDEXED",
            "extract_status": "pending",
            # ``text_quality`` is intentionally a statement about the
            # adapter state, not a claim about the source itself.  It moves
            # from ``not_examined`` to a concrete quality label only after a
            # bounded body extraction is attempted.
            "text_quality": (
                "not_examined" if extension in BODY_EXTRACTABLE_EXTENSIONS
                else "metadata_only"
            ),
            "_path": str(path),
        }

    @staticmethod
    def _body_extractable(record: Dict[str, Any]) -> bool:
        """Whether the current adapter has a bounded body extractor.

        The inventory intentionally accepts arbitrary files as metadata.  A
        separate predicate prevents search from counting an image, binary
        Office file, or unknown extension as if its body had been examined.
        """
        return str(record.get("extension") or "").lower() in BODY_EXTRACTABLE_EXTENSIONS

    @staticmethod
    def _text_quality(text: str, extension: str, *, partial: bool = False) -> str:
        """Return a conservative quality marker for an extracted preview.

        This is a retrieval-quality signal only.  It never upgrades a source
        or mathematical claim to ``verified``.  Replacement/control
        characters are treated as low-quality text so callers can route a PDF
        to OCR or ask the Owner to inspect the original.
        """
        value = str(text or "")
        if not value.strip():
            return "ocr_required" if str(extension).lower() == ".pdf" else "unavailable"
        length = max(1, len(value))
        replacement_ratio = value.count("\ufffd") / length
        control_count = sum(1 for char in value if ord(char) < 32 and char not in "\n\r\t")
        control_ratio = control_count / length
        if replacement_ratio >= 0.01 or control_ratio >= 0.02:
            return "ocr_low" if str(extension).lower() == ".pdf" else "encoding_low"
        if partial:
            return "native_text_partial"
        if str(extension).lower() == ".pdf":
            return "native_text"
        if str(extension).lower() in BODY_EXTRACTABLE_EXTENSIONS:
            return "native_text"
        return "unknown"

    @staticmethod
    def _extractability_for_records(records: Iterable[Dict[str, Any]]) -> Dict[str, Any]:
        """Summarise what the current adapter can and cannot inspect.

        Counts are based on the indexed snapshot, not a speculative full-text
        scan.  ``ocr_candidates`` therefore means PDFs that *may* need OCR;
        it is not an assertion that every PDF is image-only.
        """
        rows = list(records)
        total = len(rows)
        body_capable = [row for row in rows if str(row.get("extension") or "").lower() in BODY_EXTRACTABLE_EXTENSIONS]
        legacy = [row for row in rows if str(row.get("extension") or "").lower() in LEGACY_OFFICE_EXTENSIONS]
        unsupported = Counter(
            str(row.get("extension") or "<none>").lower()
            for row in rows
            if (
                str(row.get("extension") or "").lower() not in ALLOWED_EXTENSIONS
                and str(row.get("extension") or "").lower() not in LEGACY_OFFICE_EXTENSIONS
            )
        )
        # A metadata-only item includes unsupported files, images/assets, and
        # legacy Office formats.  Keep the category disjoint from the body
        # extractor count so rates are easy to audit.
        metadata_only = max(0, total - len(body_capable))
        status_counts = Counter(str(row.get("extract_status") or "pending") for row in rows)
        quality_counts = Counter(str(row.get("text_quality") or "not_examined") for row in rows)
        rate = (len(body_capable) / total) if total else 0.0
        result = {
            "text_or_pdf_or_openxml": len(body_capable),
            "metadata_only": metadata_only,
            "extractability_rate": round(rate, 6),
            "extractability_rate_pct": round(rate * 100.0, 4),
            "legacy_office": len(legacy),
            "ocr_candidates": sum(1 for row in rows if str(row.get("extension") or "").lower() == ".pdf"),
            "unsupported_extensions": dict(sorted(unsupported.items())),
            # Add explicit aliases/diagnostics while keeping the six compact
            # fields above stable for UI clients and tests.
            "total_indexed": total,
            "body_extractable_extensions": sorted(BODY_EXTRACTABLE_EXTENSIONS),
            "unsupported_extension_count": int(sum(unsupported.values())),
            "status_counts": dict(sorted(status_counts.items())),
            "quality_counts": dict(sorted(quality_counts.items())),
            "coverage_basis": "indexed_metadata_snapshot; body text is examined on demand",
        }
        return result

    def inventory(self, force_refresh: bool = False, include_documents: bool = False) -> Dict[str, Any]:
        with self._lock:
            cache_fresh = self._summary_cache is not None and (
                KB_CACHE_TTL_SECONDS <= 0 or time.monotonic() - self._summary_cached_at < KB_CACHE_TTL_SECONDS
            )
            if cache_fresh and not force_refresh:
                # Body extraction is on demand and can update per-record
                # quality/status after the inventory was cached.  Refresh
                # only this small aggregate instead of rescanning the 19k-file
                # source directory, so the summary remains honest without
                # making every search expensive.
                extractability = self._extractability_for_records(self._records.values())
                cached = {
                    **self._summary_cache,
                    "extractability": extractability,
                    "extractability_rate": extractability["extractability_rate"],
                }
                if include_documents:
                    return {**cached, "documents": [self._public_record(item) for item in sorted(self._records.values(), key=lambda x: x["path_rel"])]}
                return cached
            root = self._root()
            started = _now()
            if root is None:
                empty_extractability = self._extractability_for_records([])
                result = {
                    "root_id": self.root_id, "source_status": "UNAVAILABLE",
                    "scan_started_at": started, "scan_completed_at": _now(),
                    "root_label": self.root.name if self.root.name else "本地资料包",
                    "valid_count": 0, "indexed_count": 0, "temporary_count": 0, "pending_count": 0, "rejected_count": 0,
                    "bytes": 0,
                    "facets": {"modules": {}, "extensions": {}, "years": {}, "kinds": {}},
                    "catalog_claims": [], "catalog_consistent": False,
                    "index_revision": "kb:" + hashlib.sha256(b"unavailable").hexdigest(),
                    "warnings": ["材料根目录不存在或不可读"],
                    "extractability": empty_extractability,
                    # Top-level aliases make the boundary discoverable to
                    # lightweight clients while the nested object remains
                    # the canonical contract.
                    "extractability_rate": empty_extractability["extractability_rate"],
                }
                self._records = {}
                self._summary_cache = result
                self._summary_cached_at = time.monotonic()
                return result

            records: Dict[str, Dict[str, Any]] = {}
            temporary = rejected = total_bytes = 0
            stack: List[Path] = [root]
            seen_dirs = {str(root)}
            while stack:
                current = stack.pop()
                try:
                    entries = list(os.scandir(current))
                except OSError:
                    rejected += 1
                    continue
                for entry in entries:
                    path = Path(entry.path)
                    suffix = path.suffix.lower()
                    if suffix in TEMP_SUFFIXES or any(path.name.lower().endswith(item) for item in TEMP_SUFFIXES):
                        temporary += 1
                        continue
                    if entry.is_symlink():
                        # A link is accepted only when its target remains below
                        # the configured root; links to directories are not
                        # traversed to avoid cycles.
                        rel = self._relative(path, root)
                        if rel is None:
                            rejected += 1
                            continue
                        if entry.is_dir(follow_symlinks=False):
                            continue
                    if entry.is_dir(follow_symlinks=False):
                        try:
                            resolved_dir = str(path.resolve(strict=True))
                            if resolved_dir in seen_dirs:
                                continue
                            if self._relative(path, root) is None:
                                rejected += 1
                                continue
                            seen_dirs.add(resolved_dir)
                            stack.append(path)
                        except OSError:
                            rejected += 1
                        continue
                    if not entry.is_file(follow_symlinks=True):
                        rejected += 1
                        continue
                    rel = self._relative(path, root)
                    if rel is None:
                        rejected += 1
                        continue
                    try:
                        stat = path.stat()
                        size = stat.st_size
                        # Hash small files during inventory for immediate
                        # provenance.  Large PDFs/books are hashed lazily when
                        # a user opens them; this keeps the first scan usable
                        # while never fabricating a content hash.
                        digest = _sha256(path) if size <= INLINE_HASH_LIMIT else None
                    except OSError:
                        rejected += 1
                        continue
                    total_bytes += size
                    record = self._classify(
                        rel,
                        path,
                        size,
                        digest,
                        int(getattr(stat, "st_mtime_ns", 0)),
                        int(getattr(stat, "st_ctime_ns", 0)),
                        int(getattr(stat, "st_ino", 0)),
                    )
                    records[record["doc_id"]] = record

            self._records = records
            self._text_cache.clear()
            self._text_cache_pages.clear()
            facets = {"modules": {}, "extensions": {}, "years": {}, "kinds": {}}
            for record in records.values():
                for key, value in (("modules", record["module"]), ("extensions", record["extension"]), ("kinds", record["kind"])):
                    facets[key][value] = facets[key].get(value, 0) + 1
                for year in record["years"]:
                    facets["years"][str(year)] = facets["years"].get(str(year), 0) + 1
            claims = self._catalog_claims(root)
            # The catalogue is promotional metadata, not an inventory proof.
            # Keep it visible but fail closed even if future numbers happen to match.
            revision_rows = [{k: record[k] for k in ("doc_id", "path_rel", "sha256", "hash_status", "mtime_ns", "ctime_ns", "inode", "size")} for record in records.values()]
            revision = "kb:" + hashlib.sha256(self._canonical(revision_rows).encode("utf-8")).hexdigest()
            source_status = "LOCAL_PENDING" if temporary else "LOCAL_INDEXED"
            extractability = self._extractability_for_records(records.values())
            result = {
                "root_id": self.root_id, "root_label": root.name or "本地资料包", "source_status": source_status,
                "scan_started_at": started, "scan_completed_at": _now(), "last_scan_at": _now(),
                "valid_count": len(records), "indexed_count": len(records), "temporary_count": temporary,
                "pending_count": temporary,
                "rejected_count": rejected, "bytes": total_bytes, "total_bytes": total_bytes,
                "hash_policy": {
                    "inventory": "small files only",
                    "on_demand_limit_bytes": ON_DEMAND_HASH_LIMIT,
                    "large_file_status": "DEFERRED_LARGE",
                },
                "facets": facets, "catalog_claims": claims, "catalog_consistent": False,
                "index_revision": revision,
                "warnings": ["目录说明中的宣传数字仅供参考，catalog_consistent=false；以本次只读扫描为准"] if claims else [],
                "extractability": extractability,
                # Compatibility aliases for clients that prefer a flat
                # summary.  The nested ``extractability`` object is the
                # source of truth.
                "extractability_rate": extractability["extractability_rate"],
            }
            self._summary_cache = result
            self._summary_cached_at = time.monotonic()
            if include_documents:
                return {**result, "documents": [self._public_record(item) for item in sorted(records.values(), key=lambda x: x["path_rel"])]}
            return result

    def summary(self, force_refresh: bool = False, include_documents: bool = False) -> Dict[str, Any]:
        return self.inventory(force_refresh=force_refresh, include_documents=include_documents)

    @staticmethod
    def _canonical(value: Any) -> str:
        import json
        return json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":"))

    @staticmethod
    def _public_record(record: Dict[str, Any]) -> Dict[str, Any]:
        return {key: value for key, value in record.items() if not key.startswith("_")}

    def _catalog_claims(self, root: Path) -> List[Dict[str, Any]]:
        claims: List[Dict[str, Any]] = []
        for path in root.rglob("*.html"):
            # Resolve and fence the candidate before opening it.  A sync
            # directory may contain a symlink/reparse point; reading first and
            # checking later would still disclose bytes from outside the
            # owner-provided root to this process.
            rel = self._relative(path, root)
            if not rel:
                continue
            try:
                resolved = path.resolve(strict=True)
                resolved.relative_to(root)
                if not resolved.is_file():
                    continue
                text = _strip_html(_decode(path.read_bytes()[:2 * 1024 * 1024]))
            except (OSError, ValueError):
                continue
            numbers = re.findall(r"\b\d[\d,]*(?:\.\d+)?\s*(?:GB|MB|个资料文件|个文件|资料文件|文件|PDF)", text, flags=re.I)
            if numbers:
                claims.append({"path_rel": rel, "promotional_numbers": numbers[:100]})
        return claims

    def _ensure_records(self) -> None:
        if not self._records:
            self.inventory()

    @staticmethod
    def _office_archive_safe(path: Path) -> bool:
        """Reject oversized/hostile Office archives before parsing them."""
        try:
            if path.stat().st_size > OFFICE_SOURCE_LIMIT:
                return False
            with zipfile.ZipFile(path) as archive:
                total = 0
                for info in archive.infolist():
                    if info.flag_bits & 0x1:  # encrypted member
                        return False
                    total += int(info.file_size)
                    if total > OFFICE_UNPACK_LIMIT:
                        return False
            return True
        except (OSError, zipfile.BadZipFile, RuntimeError):
            return False

    def _record_text(self, record: Dict[str, Any], limit: int = 12000, pdf_pages: int = 5) -> str:
        doc_id = record["doc_id"]
        ext = record["extension"]
        if doc_id in self._text_cache and (ext != ".pdf" or self._text_cache_pages.get(doc_id, 0) >= max(1, int(pdf_pages))):
            return self._text_cache[doc_id]
        path = Path(record["_path"])
        text = ""
        extraction_blocked = False
        partial_page_error = False
        try:
            if ext in TEXT_EXTENSIONS:
                text = _decode(path.read_bytes()[: min(max(limit * 4, 65536), INLINE_TEXT_LIMIT)])
                if ext in {".html", ".htm"}:
                    text = _strip_html(text)
            elif ext == ".docx":
                if not self._office_archive_safe(path):
                    extraction_blocked = True
                else:
                    from docx import Document
                    text = "\n".join(paragraph.text for paragraph in Document(str(path)).paragraphs)
            elif ext == ".pptx":
                if not self._office_archive_safe(path):
                    extraction_blocked = True
                else:
                    from pptx import Presentation
                    slides = Presentation(str(path)).slides
                    text = "\n".join(shape.text for slide in slides for shape in slide.shapes if hasattr(shape, "text"))
            elif ext == ".xlsx":
                if not self._office_archive_safe(path):
                    extraction_blocked = True
                else:
                    from openpyxl import load_workbook
                    workbook = load_workbook(str(path), read_only=True, data_only=True)
                    rows: List[str] = []
                    for sheet in workbook.worksheets:
                        rows.append(sheet.title)
                        for row in sheet.iter_rows(max_row=100):
                            rows.append(" ".join(str(cell.value or "") for cell in row))
                    text = "\n".join(rows)
            elif ext == ".pdf":
                try:
                    from pypdf import PdfReader
                    reader = PdfReader(str(path))
                    pages = max(1, int(pdf_pages))
                    page_texts = []
                    for page in reader.pages[:pages]:
                        try:
                            page_texts.append(page.extract_text() or "")
                        except Exception:
                            # A damaged page must not erase usable text from
                            # earlier pages; preserve a partial preview and
                            # let the status/warning signal the limitation.
                            page_texts.append("")
                            partial_page_error = True
                    text = "\n".join(page_texts)
                    record["preview_pages"] = min(pages, len(reader.pages))
                    record["page_count"] = len(reader.pages)
                except Exception:
                    text = ""
        except Exception:
            text = ""
        text = re.sub(r"\s+", " ", text).strip()[:limit]
        self._text_cache[doc_id] = text
        self._text_cache_pages[doc_id] = max(1, int(pdf_pages)) if ext == ".pdf" else 0
        if extraction_blocked:
            record["extract_status"] = "PREVIEW_UNAVAILABLE_LARGE"
            record["text_quality"] = "unavailable_large"
        elif partial_page_error and text:
            record["extract_status"] = "TEXT_PARTIAL"
            record["text_quality"] = self._text_quality(text, ext, partial=True)
        elif text:
            record["extract_status"] = "TEXT_EXTRACTED"
            record["text_quality"] = self._text_quality(text, ext)
        elif ext == ".pdf":
            record["extract_status"] = "OCR_REQUIRED"
            record["text_quality"] = "ocr_required"
        else:
            record["extract_status"] = "PREVIEW_UNAVAILABLE"
            record["text_quality"] = "unavailable"
        return text

    def _ensure_hash(self, record: Dict[str, Any]) -> None:
        """Compute a deferred content hash only after an explicit document use."""
        if record.get("hash_status") != "DEFERRED":
            return
        size = int(record.get("size_bytes", record.get("size", 0)) or 0)
        if ON_DEMAND_HASH_LIMIT and size > ON_DEMAND_HASH_LIMIT:
            record["hash_status"] = "DEFERRED_LARGE"
            record["hash_deferred_reason"] = f"file exceeds on-demand limit ({ON_DEMAND_HASH_LIMIT} bytes)"
            return
        try:
            record["sha256"] = "sha256:" + _sha256(Path(record["_path"]))
            record["hash_status"] = "INLINE_ON_DEMAND"
        except OSError:
            record["hash_status"] = "UNAVAILABLE"

    @staticmethod
    def _record_is_current(record: Dict[str, Any]) -> bool:
        """Check the cheap metadata fence before exposing a cached source.

        The inventory is intentionally cached between searches.  A sync client
        can nevertheless replace a file in place, so document/file reads must
        fail closed when size, timestamps, or inode no longer match the indexed
        snapshot.  A caller can then request ``refresh=true`` explicitly.
        """
        try:
            stat = Path(record["_path"]).stat()
        except (KeyError, OSError):
            return False
        expected = (
            int(record.get("size_bytes", record.get("size", -1)) or -1),
            int(record.get("mtime_ns", -1) or -1),
            int(record.get("ctime_ns", -1) or -1),
            int(record.get("inode", -1) or -1),
        )
        actual = (
            int(getattr(stat, "st_size", -2)),
            int(getattr(stat, "st_mtime_ns", -2)),
            int(getattr(stat, "st_ctime_ns", -2)),
            int(getattr(stat, "st_ino", -2)),
        )
        # Older records may not carry ctime/inode.  Size+mtime remains the
        # compatibility fence for those records; newly scanned records use all
        # four fields in the revision.
        if record.get("ctime_ns") in (None, 0) and record.get("inode") in (None, 0):
            return actual[:2] == expected[:2]
        return actual == expected

    def _mark_stale(self, record: Dict[str, Any]) -> Dict[str, Any]:
        record["source_status"] = "SOURCE_CHANGED"
        record["extract_status"] = "REINDEX_REQUIRED"
        record["hash_status"] = record.get("hash_status") or "DEFERRED"
        record["warning"] = "源文件自索引后发生变化；请先刷新资料库快照"
        return record

    def search(self, query: str, module: Optional[str] = None, kind: Optional[str] = None,
               year: Optional[Any] = None, extension: Optional[str] = None,
               top_k: int = 8, with_preview: bool = False) -> Dict[str, Any]:
        self._ensure_records()
        try:
            top_k = int(top_k)
        except (TypeError, ValueError):
            top_k = 8
        top_k = max(1, min(20, top_k))
        query = str(query or "").strip()
        tokens = _tokens(query)
        query_years = {int(value) for value in YEAR_RE.findall(query) if 1990 <= int(value) <= datetime.now(timezone.utc).year + 1}
        query_problem_codes = {value.upper() for value in re.findall(r"\b([A-E])\s*题?\b", query, flags=re.I)}
        requested_years = {int(value) for value in YEAR_RE.findall(str(year))} if year is not None else set()
        ext_filter = str(extension or "").lower()
        if ext_filter and not ext_filter.startswith("."):
            ext_filter = "." + ext_filter
        metadata_candidates = []
        all_filtered = []
        for record in self._records.values():
            if module and str(module).lower() not in record["module"].lower():
                continue
            if kind and str(kind).lower() != record["kind"].lower():
                continue
            if ext_filter and ext_filter != record["extension"]:
                continue
            if requested_years and not requested_years.intersection(record["years"]):
                continue
            all_filtered.append(record)
            searchable = " ".join([record["path_rel"], record["title"], record["module"], " ".join(record["tags"])])
            haystack = searchable.lower()
            score = (4.0 if query and query.lower() in haystack else 0.0)
            for token in tokens:
                if token in haystack:
                    score += 1.0
                    if token in record["title"].lower():
                        score += 2.0
                    elif token in record["path_rel"].lower():
                        score += 1.0
            if query_years.intersection(record["years"]):
                score += 5.0 * len(query_years.intersection(record["years"]))
            if query_problem_codes and any(re.search(rf"\b{re.escape(code)}(?:题|[-_])?\b", record["title"], flags=re.I) for code in query_problem_codes):
                score += 3.0
            metadata_candidates.append((score, record))

        # Content extraction is deliberately bounded.  A first search must be
        # responsive even while the user directory is several gigabytes and
        # still receiving downloads.  Metadata hits are inspected first; a
        # small second pass over text-friendly files catches terms that only
        # occur in document bodies.  The response warns when that pass is
        # capped so a caller does not mistake lexical retrieval for exhaustive
        # semantic coverage.
        metadata_candidates.sort(key=lambda item: (-item[0], item[1]["path_rel"]))
        metadata_candidate_count = len(metadata_candidates)
        positive_count = sum(1 for score, _ in metadata_candidates if score > 0)
        metadata_match_ids = {
            record["doc_id"] for score, record in metadata_candidates if score > 0
        }
        if not query:
            # An empty query is an explicit metadata browse.  Every returned
            # row is therefore metadata-derived even though its score is 0.
            metadata_match_ids = {record["doc_id"] for _, record in metadata_candidates}
        # When the filename/path already answers the query, inspect only the
        # handful that will be shown.  Body extraction is reserved for the
        # harder body-only search case and remains bounded.
        inspect_limit = min(max(top_k, 8) if positive_count >= top_k else 60, len(metadata_candidates))
        inspected_ids = set()
        body_examined_ids = set()
        body_matched_ids = set()
        preview_extract_limit = min(3, top_k)
        scored = []
        for position, (base_score, record) in enumerate(metadata_candidates[:inspect_limit]):
            inspected_ids.add(record["doc_id"])
            # Empty queries are a catalog browse, not a request to parse every
            # document.  For metadata hits use the title as a safe snippet;
            # parse only the visible top rows when a body preview is useful.
            should_extract = (
                bool(query)
                and self._body_extractable(record)
                and (base_score <= 0 or position < preview_extract_limit)
            )
            body = ""
            if should_extract and len(body_examined_ids) < 12:
                body_examined_ids.add(record["doc_id"])
                body = self._record_text(record, limit=12000, pdf_pages=1)
            haystack = body.lower()
            score = base_score
            body_match = False
            if query:
                for token in tokens:
                    if token in haystack:
                        score += 1.5
                        body_match = True
            if body_match:
                body_matched_ids.add(record["doc_id"])
            if score > 0 or not query:
                scored.append((score, record, body))

        # If metadata produced fewer useful candidates, inspect a bounded set
        # of remaining small/text records for body-only matches.
        if query and len(scored) < top_k:
            for record in all_filtered:
                if record["doc_id"] in inspected_ids or len(inspected_ids) >= inspect_limit:
                    continue
                if not self._body_extractable(record):
                    continue
                if record["size_bytes"] > INLINE_TEXT_LIMIT and record["extension"] not in {".pdf", ".docx", ".pptx", ".xlsx"}:
                    continue
                inspected_ids.add(record["doc_id"])
                if len(body_examined_ids) >= 12:
                    break
                body_examined_ids.add(record["doc_id"])
                body = self._record_text(record, limit=12000, pdf_pages=1)
                haystack = body.lower()
                score = 0.0
                body_match = False
                for token in tokens:
                    if token in haystack:
                        score += 1.5
                        body_match = True
                if body_match:
                    body_matched_ids.add(record["doc_id"])
                if score > 0:
                    scored.append((score, record, body))

        results = []
        stale_count = 0
        duplicate_count = 0
        seen_result_ids = set()
        for score, record, body in scored:
            doc_id = record["doc_id"]
            if doc_id in seen_result_ids:
                duplicate_count += 1
                continue
            seen_result_ids.add(doc_id)
            if not self._record_is_current(record):
                stale_count += 1
                continue
            result = self._public_record(record)
            result["citation_ref"] = "kbdoc:" + doc_id
            result["score"] = round(float(score), 3)
            result["snippet"] = (body[:320] or f"文件命中：{record['title']}")[:320]
            body_examined = doc_id in body_examined_ids
            body_matched = doc_id in body_matched_ids
            metadata_matched = doc_id in metadata_match_ids
            if body_matched and metadata_matched:
                result["match_source"] = "metadata+body"
            elif body_matched:
                result["match_source"] = "body"
            else:
                result["match_source"] = "metadata"
            result["body_examined"] = body_examined
            result["body_matched"] = body_matched
            if with_preview:
                result["preview"] = body[:2400]
            results.append((score, result))
        results.sort(key=lambda item: (-item[0], item[1]["path_rel"]))
        warnings = []
        if len(metadata_candidates) > inspect_limit:
            warnings.append(f"正文抽取已限流：最多检查 {inspect_limit} 个候选文件；结果仍可按文件名/目录继续检索")
        if with_preview and len(body_examined_ids) < len(results):
            warnings.append("预览为受限短片段；扫描型 PDF 可能需要打开原文件或后续 OCR")
        if stale_count:
            warnings.append(f"已排除 {stale_count} 个索引后发生变化的文件；请刷新资料库快照")
        returned = results[:top_k]
        returned_rows = [item[1] for item in returned]
        metadata_only_results = sum(1 for row in returned_rows if not row.get("body_matched"))
        body_examined_count = len(body_examined_ids)
        body_matched_count = len(body_matched_ids)
        if not query:
            ranking_mode = "metadata_browse"
        elif positive_count:
            ranking_mode = "metadata_then_bounded_body_lexical"
        else:
            ranking_mode = "bounded_body_fallback_lexical"
        truncated = bool(len(results) > top_k or len(metadata_candidates) > inspect_limit)
        search_extractability = {
            "metadata_candidates": metadata_candidate_count,
            "metadata_matches": positive_count if query else metadata_candidate_count,
            "body_examined": body_examined_count,
            "body_matched": body_matched_count,
            "metadata_only_results": metadata_only_results,
            "body_match_rate": round((body_matched_count / body_examined_count) if body_examined_count else 0.0, 6),
            "result_body_coverage": round(
                (sum(1 for row in returned_rows if row.get("body_matched")) / len(returned_rows))
                if returned_rows else 0.0,
                6,
            ),
            "basis": "bounded lexical body inspection; not exhaustive semantic retrieval",
        }
        deduplication = {
            "mode": "doc_id",
            "input_results": len(scored),
            "unique_results": len(results),
            "duplicates_removed": duplicate_count,
            "exact_content_dedup": False,
            "note": "doc_id is path-derived; exact/near-content deduplication is not enabled in this slice",
        }
        retrieval_boundary = {
            "candidate_scope": "indexed_records_after_filters",
            "metadata_candidates": metadata_candidate_count,
            "metadata_match_candidates": positive_count if query else metadata_candidate_count,
            "metadata_inspect_limit": inspect_limit,
            "body_examined_limit": 12,
            "body_examined": body_examined_count,
            "pdf_page_limit": 1,
            "preview_char_limit": 2400 if with_preview else 0,
            "top_k": top_k,
            "exhaustive": False,
            "truncated": truncated,
        }
        # ``total`` remains the number of scored matches for compatibility;
        # the explicit candidate count prevents a UI from implying that a
        # bounded scan was exhaustive.
        candidate_total = positive_count if query else len(all_filtered)
        return {
            "query": query,
            "index_revision": self.summary()["index_revision"],
            "results": returned_rows,
            "total": len(results),
            "total_candidates": candidate_total,
            "returned_count": len(returned),
            "truncated": truncated,
            "warnings": warnings,
            # New auditable retrieval-boundary fields.  Keep the scalar names
            # flat for simple clients and mirror them in ``extractability`` so
            # richer clients can consume one cohesive object.
            "metadata_candidates": metadata_candidate_count,
            "metadata_matches": positive_count if query else metadata_candidate_count,
            "body_examined": body_examined_count,
            "body_matched": body_matched_count,
            "metadata_only_results": metadata_only_results,
            "deduplication": deduplication,
            "ranking_mode": ranking_mode,
            "extractability": search_extractability,
            "retrieval_boundary": retrieval_boundary,
        }

    def document(self, doc_id: str, include_preview: bool = True) -> Optional[Dict[str, Any]]:
        self._ensure_records()
        record = self._records.get(str(doc_id))
        if not record or not DOC_ID_RE.fullmatch(str(doc_id)):
            return None
        if not self._record_is_current(record):
            return self._public_record(self._mark_stale(record))
        self._ensure_hash(record)
        result = self._public_record(record)
        result["citation_ref"] = "kbdoc:" + record["doc_id"]
        if include_preview:
            preview = self._record_text(record, limit=6000)
            result["preview"] = preview
            if record["extension"] == ".pdf" and len(preview.strip()) < 40:
                result["extract_status"] = "OCR_REQUIRED"
            else:
                result["extract_status"] = "TEXT_EXTRACTED" if preview else "PREVIEW_UNAVAILABLE"
        return result

    def context(self, query: str, module: Optional[str] = None, kind: Optional[str] = None,
                year: Optional[Any] = None, extension: Optional[str] = None,
                top_k: int = 6) -> Dict[str, Any]:
        """Return a prompt-sized, source-linked context packet for an Agent.

        This is deliberately a projection of ``search`` rather than a second
        index.  It keeps the model-facing contract small and makes every item
        carry the same document-level citation and index revision shown in the
        Owner UI.  It is advice/evidence material, never an automatic claim.
        """
        result = self.search(
            query,
            module=module,
            kind=kind,
            year=year,
            extension=extension,
            top_k=min(12, max(1, int(top_k or 6))),
            with_preview=True,
        )
        items = []
        for row in result.get("results", []):
            items.append({
                key: row[key]
                for key in (
                    "doc_id", "citation_ref", "title", "path_rel", "module", "kind",
                    "years", "snippet", "preview", "extract_status", "text_quality", "hash_status",
                    "match_source", "body_examined", "body_matched",
                )
                if key in row
            })
        summary = self.summary()
        return {
            "query": result.get("query", str(query or "")),
            "index_revision": result.get("index_revision"),
            "source_status": summary.get("source_status", "UNAVAILABLE"),
            "items": items,
            "returned_count": len(items),
            "total_candidates": result.get("total_candidates", result.get("total", len(items))),
            "truncated": bool(result.get("truncated")),
            "warnings": list(result.get("warnings") or []),
            "metadata_candidates": result.get("metadata_candidates", 0),
            "metadata_matches": result.get("metadata_matches", 0),
            "body_examined": result.get("body_examined", 0),
            "body_matched": result.get("body_matched", 0),
            "metadata_only_results": result.get("metadata_only_results", 0),
            "deduplication": result.get("deduplication", {}),
            "ranking_mode": result.get("ranking_mode", "unknown"),
            "extractability": result.get("extractability", {}),
            "retrieval_boundary": result.get("retrieval_boundary", {}),
            "usage_note": "资料建议/证据线索；引用前需核对原文件、题面、页码与独立验证。",
        }

    def file_path(self, doc_id: str) -> Optional[Path]:
        self._ensure_records()
        if not DOC_ID_RE.fullmatch(str(doc_id)):
            return None
        record = self._records.get(str(doc_id))
        if not record or record["extension"] not in ALLOWED_EXTENSIONS:
            return None
        if not self._record_is_current(record):
            self._mark_stale(record)
            return None
        root = self._root()
        if root is None:
            return None
        path = Path(record["_path"])
        if self._relative(path, root) is None:
            return None
        return path
